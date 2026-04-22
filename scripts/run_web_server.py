#!/usr/bin/env python3
from __future__ import annotations

import argparse
import asyncio
import json
from contextlib import suppress
from pathlib import Path
from typing import Any

from aiohttp import WSMsgType, web

PUSH_INTERVAL_S = 1.0


class SummaryBroadcaster:
    def __init__(self, data_path: Path, interval_s: float = PUSH_INTERVAL_S) -> None:
        self._data_path = data_path
        self._interval_s = interval_s
        self._clients: set[web.WebSocketResponse] = set()
        self._cached_mtime_ns: int = -2
        self._cached_payload: str = self._encode_payload({
            "overview": {
                "generated_at": None,
                "generated_epoch": 0,
            },
            "sessions": [],
        })

    @property
    def client_count(self) -> int:
        return len(self._clients)

    def _encode_payload(self, data: dict[str, Any]) -> str:
        return json.dumps({"type": "session_summary", "data": data}, separators=(",", ":"))

    def _read_summary(self) -> dict[str, Any]:
        fallback = {
            "overview": {
                "generated_at": None,
                "generated_epoch": 0,
            },
            "sessions": [],
        }
        if not self._data_path.exists():
            return fallback
        try:
            with self._data_path.open("r", encoding="utf-8") as handle:
                payload = json.load(handle)
        except (OSError, json.JSONDecodeError):
            return fallback

        return payload if isinstance(payload, dict) else fallback

    def refresh_payload(self) -> str:
        try:
            mtime_ns = self._data_path.stat().st_mtime_ns
        except OSError:
            mtime_ns = -1

        if mtime_ns != self._cached_mtime_ns:
            self._cached_payload = self._encode_payload(self._read_summary())
            self._cached_mtime_ns = mtime_ns
        return self._cached_payload

    async def register(self, ws: web.WebSocketResponse) -> None:
        self._clients.add(ws)
        await ws.send_str(self.refresh_payload())

    def unregister(self, ws: web.WebSocketResponse) -> None:
        self._clients.discard(ws)

    async def _send(self, ws: web.WebSocketResponse, payload: str) -> None:
        if ws.closed:
            self._clients.discard(ws)
            return
        try:
            await ws.send_str(payload)
        except ConnectionResetError:
            self._clients.discard(ws)
        except RuntimeError:
            self._clients.discard(ws)

    async def run(self) -> None:
        while True:
            payload = self.refresh_payload()
            if self._clients:
                await asyncio.gather(
                    *(self._send(ws, payload) for ws in tuple(self._clients)),
                    return_exceptions=True,
                )
            await asyncio.sleep(self._interval_s)


async def index_handler(request: web.Request) -> web.FileResponse:
    web_dir: Path = request.app["web_dir"]
    return web.FileResponse(web_dir / "index.html")


async def websocket_handler(request: web.Request) -> web.WebSocketResponse:
    broadcaster: SummaryBroadcaster = request.app["broadcaster"]
    ws = web.WebSocketResponse(heartbeat=30)
    await ws.prepare(request)
    await broadcaster.register(ws)

    try:
        async for msg in ws:
            if msg.type == WSMsgType.TEXT and msg.data.strip().lower() == "ping":
                await ws.send_str('{"type":"pong"}')
            elif msg.type in (WSMsgType.CLOSE, WSMsgType.CLOSED, WSMsgType.ERROR):
                break
    finally:
        broadcaster.unregister(ws)

    return ws


async def health_handler(request: web.Request) -> web.Response:
    broadcaster: SummaryBroadcaster = request.app["broadcaster"]
    return web.json_response({"ok": True, "ws_clients": broadcaster.client_count})


async def export_pptx_handler(request: web.Request) -> web.Response:
    """Generate and serve a PPTX with architecture diagrams."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError as exc:
        return web.json_response({"error": "python-pptx not installed"}, status=500)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    ARCHITECTURES = {
        "single_agent_no_rag": {"name": "Single Agent (No RAG)", "nodes": ["Input", "Single Agent", "Output"]},
        "classical_rag": {"name": "Classical RAG", "nodes": ["Input", "Retriever", "Context Packer", "Single Agent", "Output"]},
        "hybrid_fast": {"name": "Hybrid Fast", "nodes": ["Input", "Router", "Fast Path", "RAG Path", "Output"]},
        "non_rag_multi_agent": {"name": "Multi-Agent (No RAG)", "nodes": ["Input", "Planner", "Diagnoser", "Tutor", "Critic", "Aggregator", "Output"]},
        "agentic_rag": {"name": "Agentic RAG", "nodes": ["Input", "Retriever", "Planner", "Tutor", "Critic", "Aggregator", "Output"]},
    }

    for key, arch in ARCHITECTURES.items():
        slide = prs.slides.add_slide(blank_layout)
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = arch["name"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        p.alignment = PP_ALIGN.LEFT

        # Nodes as boxes
        nodes = arch["nodes"]
        total = len(nodes)
        start_x = 0.5
        end_x = 12.5
        spacing = (end_x - start_x) / max(1, total - 1)
        for i, node in enumerate(nodes):
            x = start_x + spacing * i
            box = slide.shapes.add_shape(1, Inches(x), Inches(3.2), Inches(1.4), Inches(0.8))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0x4E, 0xCD, 0xC4) if i % 2 == 1 else RGBColor(0xF5, 0xD5, 0x47)
            box.line.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            box.line.width = Pt(2)
            tf2 = box.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = node
            p2.font.size = Pt(14)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            p2.alignment = PP_ALIGN.CENTER
            tf2.word_wrap = True

        # Connector arrows
        for i in range(total - 1):
            x1 = start_x + spacing * i + 1.4
            y1 = 3.2 + 0.4
            x2 = start_x + spacing * (i + 1)
            connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y1))
            connector.line.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            connector.line.width = Pt(2)

    from io import BytesIO
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return web.Response(body=buf.read(), content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": "attachment; filename=architecture-diagrams.pptx"})


async def on_startup(app: web.Application) -> None:
    broadcaster: SummaryBroadcaster = app["broadcaster"]
    app["broadcast_task"] = asyncio.create_task(broadcaster.run())


async def on_cleanup(app: web.Application) -> None:
    task: asyncio.Task[Any] | None = app.get("broadcast_task")
    if task is None:
        return
    task.cancel()
    with suppress(asyncio.CancelledError):
        await task


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Serve dashboard with websocket updates")
    parser.add_argument("--port", type=int, default=8080)
    parser.add_argument("--host", default="0.0.0.0")
    parser.add_argument("--web-dir", default="web")
    return parser.parse_args()


def build_app(web_dir: Path) -> web.Application:
    app = web.Application()
    data_path = web_dir / "data" / "session_summary.json"
    app["web_dir"] = web_dir
    app["broadcaster"] = SummaryBroadcaster(data_path=data_path)
    app.router.add_get("/ws", websocket_handler)
    app.router.add_get("/healthz", health_handler)
    app.router.add_get("/export/pptx", export_pptx_handler)
    app.router.add_get("/", index_handler)
    app.router.add_static("/", path=str(web_dir), show_index=False)
    app.on_startup.append(on_startup)
    app.on_cleanup.append(on_cleanup)
    return app


def main() -> None:
    args = parse_args()
    web_dir = Path(args.web_dir).expanduser().resolve()
    if not web_dir.exists():
        raise SystemExit(f"Web directory does not exist: {web_dir}")

    app = build_app(web_dir)
    print(f"Serving {web_dir} at http://{args.host}:{args.port} with websocket /ws (1s push)")
    web.run_app(app, host=args.host, port=args.port)


if __name__ == "__main__":
    main()
